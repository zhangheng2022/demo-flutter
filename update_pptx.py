# -*- coding: utf-8 -*-
"""
根据 demo_flutter 项目完善 AI 实战落地汇报模板
"""
from pptx import Presentation

REPLACEMENTS = [
    # (原文本, 新文本) - 按从长到短排序，优先匹配长串
    ('【部门名称】', '技术研发部'),
    ('[填写本部门AI应用的核心场景名称]', 'AI 辅助 Flutter 移动应用开发'),
    ('[简要描述业务痛点]', '传统手动开发周期长；多模块（拍照、扫码、手机轮廓拍摄）实现复杂；需快速验证产品原型'),
    ('简要描述业务痛点', '传统手动开发周期长；多模块（拍照、扫码、手机轮廓拍摄）实现复杂；需快速验证产品原型'),
    ('[如：ChatGPT、Claude、文心一言、Midjourney等]', 'Cursor、GitHub Copilot、ChatGPT（架构设计、代码生成、问题排查）'),
    ('如：ChatGPT、Claude、文心一言、Midjourney等', 'Cursor、GitHub Copilot、ChatGPT（架构设计、代码生成、问题排查）'),
    ('[具体描述]', '开发 demo_flutter 相机应用：包含普通拍照、手机多角度轮廓拍照、条形码扫描、相册管理、扫码结果列表等功能模块'),
    ('[步骤1 → 步骤2 → 步骤3]', '需求拆分 → AI 生成框架与 Provider 设计 → 逐模块实现（相机、扫描、相册）→ 联调与优化'),
    ('步骤1 → 步骤2 → 步骤3', '需求拆分 → AI 生成框架与 Provider 设计 → 逐模块实现（相机、扫描、相册）→ 联调与优化'),
    ('[人员名单]', '研发团队'),
    ('人员名单', '研发团队'),
    ('[在此粘贴\n实际使用截图\n或效果演示图]', '（建议插入：应用首页截图、扫码界面、手机轮廓拍照界面）'),
    ('[（建议插入截图或演示图）', '（建议插入：应用首页、扫码、手机轮廓拍照界面截图）'),
    ('在此粘贴', '（建议插入：应用首页、扫码、手机轮廓拍照界面截图）'),
    ('XX%', '40-60%'),
    ('[具体数值]', '待统计（可填写实际节省的小时数或金额）'),
    ('具体数值', '待统计'),
    ('[计划开展的下一个AI应用场景]', 'AI 辅助单元测试与集成测试生成；AI 辅助 UI 设计稿转 Flutter 代码'),
    ('计划开展的下一个AI应用场景', 'AI 辅助单元测试与集成测试生成；AI 辅助 UI 设计稿转 Flutter 代码'),
    ('[预期达成的效果]', '提升测试覆盖率、缩短 UI 开发周期'),
    ('预期达成的效果', '提升测试覆盖率、缩短 UI 开发周期'),
    ('[预计完成时间]', '2026年Q2'),
    ('预计完成时间', '2026年Q2'),
]


def replace_in_shape(shape):
    """在 shape 的 run 级别替换文本"""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in REPLACEMENTS:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        break
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for old, new in REPLACEMENTS:
                    if old in cell.text:
                        cell.text = cell.text.replace(old, new)


def main():
    prs = Presentation('AI实战落地汇报模板.pptx')
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape)
    prs.save('AI实战落地汇报模板.pptx')
    print('已成功更新 AI实战落地汇报模板.pptx')


if __name__ == '__main__':
    main()
